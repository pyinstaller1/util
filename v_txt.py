



import fitz  # PyMuPDF
import pandas as pd
from pptx import Presentation
from docx import Document


def get_pdf(path):
    doc = fitz.open(path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text


def get_txt(path, encoding='utf-8'):
    with open(path, 'r', encoding=encoding) as f:
        return f.read()


def get_csv(path, encoding='utf-8'):
    df = pd.read_csv(path, encoding=encoding)
    df = df.fillna("").astype(str)  # NaN 제거, 문자열로 변환
    text = "\n".join(["\t".join(row) for row in df.values])
    return text

def get_xlsx(path, sheet_name=0):
    df = pd.read_excel(path, sheet_name=sheet_name) # 1번째 sheet
    text = df.fillna("").astype(str).values     # NaN 제거, 문자열로 변환
    combined_text = "\n".join(["\t".join(row) for row in text])   # 2차원 배열을 문자열로 변환
    return combined_text


def get_pptx(path):
    prs = Presentation(path)
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)

    return "\n".join(text_runs)

def get_docx(file_path):
    doc = Document(file_path)
    text = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(text)   # 리스트를 "\n"로 연결



print(7)
text = ''
text += get_pdf(".\건강보험자료.pdf")
text += get_txt(".\건강.txt")


print(text)

print(8)










